from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Employee Salary Prediction"
subtitle.text = "Name: Hiren Chauhan\nEmail: hinesh13@gmail.com\nAICTE Reg. no: STU67e4d5f773a9f1743050231\nBatch: AICTE B2_AI - (2025-26)"

# Slide 2: Introduction & Problem Statement
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Introduction & Problem Statement"
content.text = (
    "- Project Goal: To develop a machine learning model that predicts whether an individual's annual income is greater than $50,000.\n"
    "- Problem: Predicting income level is a common benchmark problem in machine learning and has real-world applications in areas like financial planning and targeted marketing.\n"
    "- Solution: We have built a web application that uses a trained model to make instant predictions based on user-provided demographic and employment data."
)

# Slide 3: Dataset Overview
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Dataset Overview"
content.text = (
    '- Dataset: The project uses the "adult_3.csv" dataset, a well-known dataset from the UCI Machine Learning Repository.\n'
    "- Features: The dataset contains 15 features, including: age, workclass, education, marital-status, occupation, relationship, race, gender, capital-gain, capital-loss, hours-per-week, native-country\n"
    "- Target Variable: income (either '<=50K' or '>50K')"
)

# Slide 4: Data Preprocessing
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Data Preprocessing"
content.text = (
    "- Handling Missing Values: The dataset contained missing values represented by '?'. These were replaced with the mode (most frequent value) of their respective columns.\n"
    "- Feature Engineering: The 'fnlwgt' (final weight) column was removed as it is not relevant for individual income prediction.\n"
    "- Categorical Data Encoding: All categorical features (like workclass, occupation, gender, etc.) were converted into numerical format using LabelEncoder from scikit-learn. This is necessary for the machine learning model to process the data."
)

# Slide 5: Model Training & Evaluation
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Model Training & Evaluation"
content.text = (
    "- Algorithm: A RandomForestClassifier was chosen for this task. It is a powerful and robust ensemble learning method that generally provides high accuracy.\n"
    "- Training Process:\n"
    "  1. The data was split into a training set (80%) and a testing set (20%).\n"
    "  2. The RandomForestClassifier model was trained on the training data.\n"
    "- Evaluation:\n"
    "  - The model's performance was evaluated on the unseen test data.\n"
    "  - The primary metric used was accuracy, which measures the percentage of correct predictions.\n"
    "  - Our model achieved an accuracy of approximately 85.51%."
)

# Slide 6: Application Demo
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Application Demo"
content.text = (
    "- Technology: The prediction model is deployed in an interactive web application built with Streamlit.\n"
    "- User Interface:\n"

    "  - The application features a simple sidebar where users can input employee information (age, workclass, education, etc.).\n"
    "  - The inputs are then fed to the trained model.\n"
    "- Output:\n"
    "  - The application displays the predicted income category (>50K or <=50K).\n"
    "  - It also shows the confidence level of the prediction.\n"
    "  - A preview of the dataset is also available for reference."
)

# Slide 7: Conclusion
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Conclusion"
content.text = (
    "- We have successfully built and deployed a machine learning model to predict employee income with a reasonable accuracy of over 85%.\n"
    "- The project demonstrates a complete machine learning workflow: from data cleaning and preprocessing to model training, evaluation, and deployment in a user-friendly web application.\n"
    "- Future Work:\n"
    "  - Experiment with other classification algorithms (e.g., Gradient Boosting, XGBoost) to potentially improve accuracy.\n"
    "  - Perform more in-depth feature engineering.\n"
    "  - Deploy the application to a cloud service for public access."
)

# Slide 8: Thank You & Questions
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Thank You!"

# Save the presentation
prs.save("Employee_Salary_Prediction.pptx")

print("Presentation 'Employee_Salary_Prediction.pptx' created successfully.")
